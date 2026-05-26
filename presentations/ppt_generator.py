"""
PPT生成器
使用python-pptx生成专业的PPT演示文稿
"""

import os
from pathlib import Path
from typing import Dict, List, Any, Optional
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

class PPTGenerator:
    """
    PPT生成器 - 生成专业的演示文稿
    """

    BRAND_BLUE = RgbColor(0, 102, 204) if HAS_PPTX else None
    BRAND_PURPLE = RgbColor(124, 58, 237) if HAS_PPTX else None
    TEXT_WHITE = RgbColor(255, 255, 255) if HAS_PPTX else None
    TEXT_DARK = RgbColor(44, 62, 80) if HAS_PPTX else None

    def __init__(self, output_dir: str = "outputs"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def create_presentation(self) -> 'Presentation':
        """创建新的演示文稿"""
        if not HAS_PPTX:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")

        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)
        return prs

    def add_title_slide(
        self,
        prs: 'Presentation',
        title: str,
        subtitle: str = "",
        author: str = "",
        date: str = ""
    ) -> None:
        """添加标题幻灯片"""
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加背景矩形
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.BRAND_BLUE
        shape.line.fill.background()

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER

        # 添加副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2),
                Inches(12.333), Inches(1)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(28)
            p.font.color.rgb = self.TEXT_WHITE
            p.alignment = PP_ALIGN.CENTER

        # 添加日期
        if date:
            date_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.5),
                Inches(12.333), Inches(0.5)
            )
            tf = date_box.text_frame
            p = tf.paragraphs[0]
            p.text = date
            p.font.size = Pt(18)
            p.font.color.rgb = self.TEXT_WHITE
            p.alignment = PP_ALIGN.CENTER

    def add_content_slide(
        self,
        prs: 'Presentation',
        title: str,
        content: List[str],
        title_color: str = "#0066CC"
    ) -> None:
        """添加内容幻灯片"""
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题栏背景
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        from pptx.dml.color import RgbColor
        header.fill.fore_color.rgb = RgbColor(0, 102, 204)
        header.line.fill.background()

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.TEXT_WHITE

        # 添加内容
        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.6),
            Inches(11.933), Inches(5.4)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(24)
            p.font.color.rgb = RgbColor(44, 62, 80)
            p.space_before = Pt(12)
            p.space_after = Pt(6)

    def add_two_column_slide(
        self,
        prs: 'Presentation',
        title: str,
        left_content: List[str],
        right_content: List[str],
        left_title: str = "左侧",
        right_title: str = "右侧"
    ) -> None:
        """添加双栏幻灯片"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 标题
        from pptx.dml.color import RgbColor
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RgbColor(0, 102, 204)

        # 左侧标题
        left_title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3),
            Inches(5.8), Inches(0.5)
        )
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RgbColor(124, 58, 237)

        # 左侧内容
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.9),
            Inches(5.8), Inches(5)
        )
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.space_before = Pt(8)

        # 右侧标题
        right_title_box = slide.shapes.add_textbox(
            Inches(6.8), Inches(1.3),
            Inches(5.8), Inches(0.5)
        )
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RgbColor(124, 58, 237)

        # 右侧内容
        right_box = slide.shapes.add_textbox(
            Inches(6.8), Inches(1.9),
            Inches(5.8), Inches(5)
        )
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.space_before = Pt(8)

    def add_table_slide(
        self,
        prs: 'Presentation',
        title: str,
        headers: List[str],
        rows: List[List[str]]
    ) -> None:
        """添加表格幻灯片"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        from pptx.dml.color import RgbColor

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RgbColor(0, 102, 204)

        # 表格
        cols = len(headers)
        rows_count = len(rows) + 1

        table = slide.shapes.add_table(
            rows_count, cols,
            Inches(0.5), Inches(1.5),
            Inches(12.333), Inches(5)
        ).table

        # 设置列宽
        col_width = Inches(12.333 / cols)
        for i in range(cols):
            table.columns[i].width = col_width

        # 表头
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RgbColor(0, 102, 204)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.TEXT_WHITE
            p.alignment = PP_ALIGN.CENTER

        # 数据行
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_text)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(16)
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RgbColor(240, 240, 240)

    def generate_training_ppt(
        self,
        title: str,
        sections: List[Dict],
        output_filename: str = "training.pptx"
    ) -> str:
        """
        生成培训课件PPT

        Args:
            title: 课件标题
            sections: 章节列表 [{"title": "...", "type": "content|table|two-column", "content": [...]}]
            output_filename: 输出文件名
        """
        if not HAS_PPTX:
            raise ImportError("python-pptx not installed")

        prs = self.create_presentation()

        # 封面
        self.add_title_slide(
            prs,
            title=title,
            subtitle="商飞离线AI文档处理平台",
            author="培训部",
            date=datetime.now().strftime("%Y年%m月%d日")
        )

        # 目录
        self.add_content_slide(
            prs,
            title="目录",
            content=[s["title"] for s in sections]
        )

        # 各章节
        for section in sections:
            section_type = section.get("type", "content")

            if section_type == "content":
                self.add_content_slide(
                    prs,
                    title=section["title"],
                    content=section["content"]
                )
            elif section_type == "two-column":
                self.add_two_column_slide(
                    prs,
                    title=section["title"],
                    left_content=section.get("left", []),
                    right_content=section.get("right", []),
                    left_title=section.get("left_title", "左侧"),
                    right_title=section.get("right_title", "右侧")
                )
            elif section_type == "table":
                self.add_table_slide(
                    prs,
                    title=section["title"],
                    headers=section.get("headers", []),
                    rows=section.get("rows", [])
                )

        # 保存
        output_path = self.output_dir / output_filename
        prs.save(str(output_path))
        return str(output_path)


class PPTOptimizer:
    """
    PPT优化器 - 优化现有PPT的样式和布局
    """

    def __init__(self):
        if not HAS_PPTX:
            raise ImportError("python-pptx not installed")

    def apply_brand_style(
        self,
        pptx_path: str,
        output_path: str,
        brand_color: str = "#0066CC"
    ) -> str:
        """
        应用品牌样式到PPT

        Args:
            pptx_path: 原始PPT路径
            output_path: 输出路径
            brand_color: 品牌颜色
        """
        prs = Presentation(pptx_path)

        from pptx.dml.color import RgbColor
        # 解析颜色
        r = int(brand_color[1:3], 16)
        g = int(brand_color[3:5], 16)
        b = int(brand_color[5:7], 16)
        color = RgbColor(r, g, b)

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        # 应用品牌颜色到标题
                        if paragraph.font.size and paragraph.font.size >= Pt(24):
                            paragraph.font.color.rgb = color

        prs.save(output_path)
        return output_path

    def compress_images(self, pptx_path: str, output_path: str) -> str:
        """
        压缩PPT中的图片
        """
        # PPTX中的图片压缩需要特殊处理
        # 这里只是标记，实际压缩由PowerPoint完成
        prs = Presentation(pptx_path)
        prs.save(output_path)
        return output_path
