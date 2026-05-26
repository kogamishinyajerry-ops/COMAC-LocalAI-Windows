from pptx import Presentation
from pathlib import Path
from parsers.base_parser import BaseParser, Document

class PPTParser(BaseParser):
    def parse(self, file_path: str) -> Document:
        prs = Presentation(file_path)
        slides_content = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            slides_content.append(f"=== Slide {i+1} ===\n" + "\n".join(slide_text))

        return Document(
            content="\n\n".join(slides_content),
            metadata={
                "slides": len(prs.slides),
            },
            format="pptx",
            pages=slides_content
        )

    def extract_text(self, file_path: str) -> str:
        return self.parse(file_path).content
